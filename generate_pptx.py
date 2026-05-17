from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors (matching the app's dark theme) ──
BG          = RGBColor(0x0B, 0x0D, 0x11)
CARD_BG     = RGBColor(0x16, 0x19, 0x22)
BORDER      = RGBColor(0x1E, 0x22, 0x30)
TEXT        = RGBColor(0xE4, 0xE5, 0xE9)
TEXT2       = RGBColor(0x9C, 0xA3, 0xAF)
TEXT_MUTED  = RGBColor(0x6B, 0x72, 0x80)
ACCENT      = RGBColor(0x63, 0x66, 0xF1)
ACCENT_LIGHT = RGBColor(0x81, 0x8C, 0xF8)
ACCENT_BG   = RGBColor(0x1A, 0x1B, 0x3A)
GREEN       = RGBColor(0x22, 0xC5, 0x5E)
GREEN_BG    = RGBColor(0x16, 0x3B, 0x19)
RED         = RGBColor(0xEF, 0x44, 0x44)
RED_BG      = RGBColor(0x3B, 0x16, 0x16)
YELLOW      = RGBColor(0xEA, 0xB3, 0x08)
YELLOW_BG   = RGBColor(0x3B, 0x2E, 0x16)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Microsoft YaHei"  # CJK-friendly fallback
MONO = "Consolas"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ═══════════════════════════════════════════
# Helper functions
# ═══════════════════════════════════════════

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_rect(slide, left, top, width, height, fill_color=CARD_BG, border_color=None):
    """Add a filled rectangle with optional border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=Pt(14),
             color=TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT):
    """Add a text box with single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=Pt(11),
                  color=TEXT2, line_spacing=Pt(18), font_name=FONT):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, fs, clr, bld) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = fs or font_size
        p.font.color.rgb = clr or color
        p.font.bold = bld or False
        p.font.name = font_name
        p.space_after = line_spacing
    return txBox

def add_badge(slide, left, top, text, bg_color=ACCENT_BG, text_color=ACCENT_LIGHT):
    """Add a small badge pill."""
    shape = add_rect(slide, left, top, Inches(2.0), Inches(0.38), fill_color=bg_color)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

def add_card_with_content(slide, left, top, width, height, title, body_lines,
                          title_color=TEXT, border_top_color=None, number=None,
                          number_color=None, number_bg=None):
    """Add a card with optional numbered icon, title and body text."""
    shape = add_rect(slide, left, top, width, height)
    if border_top_color:
        # Simulate top border with a thin rectangle
        add_rect(slide, left, top, width, Pt(3), fill_color=border_top_color)

    y_offset = top + Inches(0.3)
    if number is not None:
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.3), y_offset, Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = number_bg or ACCENT_BG
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = str(number)
        p.font.size = Pt(16)
        p.font.color.rgb = number_color or ACCENT_LIGHT
        p.font.bold = True
        p.font.name = MONO
        p.alignment = PP_ALIGN.CENTER
        title_left = left + Inches(1.0)
    else:
        title_left = left + Inches(0.3)

    add_text(slide, title_left, y_offset, width - Inches(1.3), Inches(0.4),
             title, Pt(14), title_color, bold=True)

    if body_lines:
        add_multiline(slide, left + Inches(0.3), y_offset + Inches(0.55),
                      width - Inches(0.6), height - Inches(0.9), body_lines,
                      font_size=Pt(10), color=TEXT2, line_spacing=Pt(6))


# ═══════════════════════════════════════════
# SLIDE 1: 一句话 —— AgentBench 是什么
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)

add_badge(slide, Inches(1), Inches(1.2), "v0.5.0 · MIT 开源",
          bg_color=ACCENT_BG, text_color=ACCENT_LIGHT)

add_text(slide, Inches(1), Inches(1.9), Inches(11), Inches(1.2),
         "AgentBench", Pt(72), ACCENT_LIGHT, bold=True)

add_text(slide, Inches(1), Inches(3.1), Inches(8), Inches(0.6),
         "AI Agent 评测运维平台", Pt(26), TEXT, bold=True)

add_text(slide, Inches(1), Inches(3.8), Inches(9), Inches(0.8),
         "用 AI 测试 AI —— 像测试软件一样测试你的 Agent\n回归测试、漂移检测、质量保障，一站完成",
         Pt(15), TEXT2)

meta_items = [
    "Next.js 16 + React 19 + TypeScript",
    "多提供商 AI（Anthropic + OpenAI + 自定义）",
    "72 测试 · 零 Lint · 零构建错误",
    "完全由 Claude Code 构建",
]
for i, item in enumerate(meta_items):
    add_text(slide, Inches(1 + i * 3.0), Inches(5.2), Inches(2.8), Inches(0.4),
             "●  " + item, Pt(11), TEXT_MUTED)


# ═══════════════════════════════════════════
# SLIDE 2: 三个亮点（数据法 > 模式法 > 团队法）
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)

add_text(slide, Inches(1), Inches(0.5), Inches(6), Inches(0.6),
         "三个核心亮点", Pt(36), TEXT, bold=True)
add_badge(slide, Inches(7.2), Inches(0.65), "按 BP 优先级排序",
          bg_color=ACCENT_BG, text_color=ACCENT_LIGHT)

# Card 1: 数据法
card_w = Inches(3.6)
card_h = Inches(5.5)
gap = Inches(0.3)
start_x = Inches(1)

add_card_with_content(slide, start_x, Inches(1.6), card_w, card_h,
    "数据法 — 硬指标说话",
    [
        ("用数据证明质量，而非堆砌形容词。", Pt(11), TEXT2, False),
        ("", Pt(6), TEXT2, False),
        ("72              6                0", Pt(22), GREEN, True),
        ("测试通过       评测器类型      TS/Lint 错误", Pt(9), TEXT_MUTED, False),
    ],
    title_color=GREEN, border_top_color=GREEN, number=1,
    number_color=GREEN, number_bg=GREEN_BG)

# Card 2: 模式法
add_card_with_content(slide, start_x + card_w + gap, Inches(1.6), card_w, card_h,
    "模式法 — 差异化路径",
    [
        ("市面上唯一的\"本地优先 + 多提供商 AI + 无后端\"评测平台。", Pt(11), TEXT2, False),
        ("", Pt(6), TEXT2, False),
        ("• 数据全存 localStorage，导出/导入 JSON", Pt(10), TEXT_MUTED, False),
        ("• 支持 OpenAI、Anthropic、自定义 HTTP", Pt(10), TEXT_MUTED, False),
        ("• 无需数据库、无需服务器，可自行部署", Pt(10), TEXT_MUTED, False),
    ],
    title_color=ACCENT_LIGHT, border_top_color=ACCENT, number=2,
    number_color=ACCENT_LIGHT, number_bg=ACCENT_BG)

# Card 3: 团队法
add_card_with_content(slide, start_x + (card_w + gap) * 2, Inches(1.6), card_w, card_h,
    "团队法 — AI 原生构建",
    [
        ("整个平台完全由 Claude Code 构建。AI 编码既是手段，也是作品本身。", Pt(11), TEXT2, False),
        ("", Pt(6), TEXT2, False),
        ("• 3 轮迭代 · 20+ bug 修复", Pt(10), TEXT_MUTED, False),
        ("• 18+ 文件中文化本地化", Pt(10), TEXT_MUTED, False),
        ("• CC Switch 联动 · 真实数据驱动评测", Pt(10), TEXT_MUTED, False),
    ],
    title_color=YELLOW, border_top_color=YELLOW, number=3,
    number_color=YELLOW, number_bg=YELLOW_BG)


# ═══════════════════════════════════════════
# SLIDE 3: 为什么需要 AgentBench（市场痛点）
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "为什么需要 AgentBench", Pt(36), TEXT, bold=True)
add_text(slide, Inches(1), Inches(1.1), Inches(10), Inches(0.5),
         "\"市场分析\"本是最容易被跳过的环节。但如果你在做 AI Agent，以下三个事实你绕不开。",
         Pt(13), TEXT_MUTED)

pain_points = [
    ("⚠", "Agent 质量极其脆弱", RED, RED_BG,
     "提示词调整、模型升级、微调更新，任何改动都可能无声地降低 Agent 性能。"
     "没有系统化测试，回归问题不会被发现——直到用户投诉。"),
    ("⚙", "AI Agent 缺少标准 QA", YELLOW, YELLOW_BG,
     "传统软件 QA 工具（JUnit、Selenium、Jest）无法测试 AI Agent。"
     "你只能用 AI 来测试 AI——将实际响应与预期行为进行系统化对比。"),
    ("⚡", "评测运维是下一片蓝海", ACCENT_LIGHT, ACCENT_BG,
     "在 46 天的机会分析中，Agent 回归测试需求出现了 12+ 次。"
     "随着 AI Agent 从实验走向生产，评测运维将成为每个 Agent 团队的必备基础设施。"),
]

for i, (icon, title, clr, clr_bg, desc) in enumerate(pain_points):
    left = Inches(1 + i * 3.9)
    top = Inches(1.9)
    shape = add_rect(slide, left, top, Inches(3.6), Inches(3.5))
    # Icon box
    icon_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         left + Inches(0.3), top + Inches(0.3),
                                         Inches(0.55), Inches(0.55))
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = clr_bg
    icon_shape.line.fill.background()
    tf = icon_shape.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(18)
    p.font.color.rgb = clr
    p.alignment = PP_ALIGN.CENTER
    # Title
    add_text(slide, left + Inches(1.0), top + Inches(0.3), Inches(2.3), Inches(0.4),
             title, Pt(14), TEXT, bold=True)
    # Description
    add_text(slide, left + Inches(0.3), top + Inches(1.1), Inches(3.0), Inches(2.2),
             desc, Pt(11), TEXT2)

add_text(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
         "\"用 AI 测试 AI\" —— 这不是口号，这是 Agent 时代的质量范式",
         Pt(16), ACCENT_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SLIDE 4: 产品功能全景
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "产品功能全景", Pt(36), TEXT, bold=True)
add_text(slide, Inches(1), Inches(1.1), Inches(10), Inches(0.5),
         "具象展示胜过千言万语。六大核心功能模块。", Pt(13), TEXT_MUTED)

features = [
    ("测试套件管理",
     "创建、编辑、导入（JSON/JSONL）、删除。每个用例可自定义评测器。"),
    ("多评测器引擎",
     "精确匹配 · 包含（关键短语/否定感知）· 正则 · JSON Schema · LLM 评判 · 代码测试"),
    ("真实与模拟模式",
     "直连 Agent 端点（OpenAI/Anthropic/自定义 HTTP），或运行模拟回归测试。"),
    ("版本回归对比",
     "并排版本比较，含得分差异、通过率变化、不稳定检测、可视化柱状图。"),
    ("AI 驱动分析",
     "多提供商 AI 识别回归模式并建议修复方案。无 API Key 时优雅降级至演示模式。"),
    ("本地优先 · 零后端",
     "所有数据存储于 localStorage。导出/导入 JSON。无需数据库，可自行部署。"),
]

for i, (title, desc) in enumerate(features):
    row = i // 3
    col = i % 3
    left = Inches(1 + col * 3.9)
    top = Inches(1.8 + row * 2.6)
    shape = add_rect(slide, left, top, Inches(3.6), Inches(2.2))
    # Check mark
    add_text(slide, left + Inches(0.25), top + Inches(0.2), Inches(3.0), Inches(0.4),
             "✓  " + title, Pt(12), GREEN, bold=True)
    add_text(slide, left + Inches(0.25), top + Inches(0.8), Inches(3.1), Inches(1.2),
             desc, Pt(10), TEXT2)


# ═══════════════════════════════════════════
# SLIDE 5: 数据表现（最关键的一页）
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "数据表现", Pt(36), TEXT, bold=True)
add_text(slide, Inches(1), Inches(1.1), Inches(10), Inches(0.5),
         "不需要大，但需要证明\"行得通\"。", Pt(13), TEXT_MUTED)

# Stats row
stats = [
    ("72", "测试通过", GREEN),
    ("3", "测试文件", ACCENT_LIGHT),
    ("0", "TypeScript 错误", GREEN),
    ("0", "ESLint 警告", GREEN),
]
for i, (val, label, clr) in enumerate(stats):
    left = Inches(1 + i * 3.0)
    top = Inches(1.9)
    add_rect(slide, left, top, Inches(2.6), Inches(1.8))
    add_text(slide, left, top + Inches(0.2), Inches(2.6), Inches(0.7),
             val, Pt(36), clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(1.0), Inches(2.6), Inches(0.4),
             label, Pt(11), TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Coverage details card (left)
left_card = Inches(1)
top_card = Inches(4.1)
card_w2 = Inches(5.5)
add_rect(slide, left_card, top_card, card_w2, Inches(2.8))
add_text(slide, left_card + Inches(0.3), top_card + Inches(0.2), Inches(5), Inches(0.4),
         "测试覆盖详情", Pt(13), TEXT, bold=True)

coverage = [
    ("evaluator.test.ts", "43 测试", 0.60, GREEN),
    ("utils.test.ts", "14 测试", 0.19, ACCENT),
    ("agent-adapter.test.ts", "15 测试", 0.21, ACCENT_LIGHT),
]
for i, (file, count, pct, clr) in enumerate(coverage):
    y = top_card + Inches(0.8 + i * 0.65)
    add_text(slide, left_card + Inches(0.3), y, Inches(3.0), Inches(0.3),
             file, Pt(10), TEXT2)
    add_text(slide, left_card + Inches(3.6), y, Inches(1.2), Inches(0.3),
             count, Pt(10), clr, bold=True, font_name=MONO)
    # Progress bar
    bar = add_rect(slide, left_card + Inches(0.3), y + Inches(0.3),
                    Inches(4.9) * pct, Pt(4), fill_color=clr)
    add_rect(slide, left_card + Inches(0.3), y + Inches(0.3),
             Inches(4.9), Pt(4), fill_color=BORDER)
    # Re-add fill on top
    bar_shape = add_rect(slide, left_card + Inches(0.3), y + Inches(0.3),
                         Inches(4.9) * pct, Pt(4), fill_color=clr)

# Evaluator tags card (right)
right_card = Inches(7.0)
add_rect(slide, right_card, top_card, Inches(5.3), Inches(2.8))
add_text(slide, right_card + Inches(0.3), top_card + Inches(0.2), Inches(4.5), Inches(0.4),
         "评测器体系", Pt(13), TEXT, bold=True)

tags = [
    ("包含匹配", ACCENT_BG, ACCENT_LIGHT),
    ("精确匹配", ACCENT_BG, ACCENT_LIGHT),
    ("正则", ACCENT_BG, ACCENT_LIGHT),
    ("JSON Schema", ACCENT_BG, ACCENT_LIGHT),
    ("LLM 评判", ACCENT_BG, ACCENT_LIGHT),
    ("代码测试", GREEN_BG, GREEN),
]
for i, (tag, bg, clr) in enumerate(tags):
    col = i % 3
    row = i // 3
    tx = right_card + Inches(0.3 + col * 1.7)
    ty = top_card + Inches(0.8 + row * 0.6)
    tag_shape = add_rect(slide, tx, ty, Inches(1.5), Inches(0.4), fill_color=bg)
    tf = tag_shape.text_frame
    p = tf.paragraphs[0]
    p.text = tag
    p.font.size = Pt(9)
    p.font.color.rgb = clr
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

add_text(slide, right_card + Inches(0.3), top_card + Inches(2.0), Inches(4.5), Inches(0.6),
         "从文本匹配到语义理解，从结构验证到代码质量——正交覆盖 Agent 评测的所有维度。",
         Pt(10), TEXT_MUTED)


# ═══════════════════════════════════════════
# SLIDE 6: 竞品分析（比较表）
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "竞品分析", Pt(36), TEXT, bold=True)
add_text(slide, Inches(1), Inches(1.1), Inches(10), Inches(0.5),
         "没有直接对手，但传统工具无法解决 AI Agent 评测问题。", Pt(13), TEXT_MUTED)

# Build comparison table
table_data = [
    ["维度", "AgentBench", "传统 QA 工具", "手工评测"],
    ["AI Agent 评测", "✓ 专为此设计", "✗ 不适用", "◑ 低效、不可扩展"],
    ["语义理解", "✓ LLM 评判 + 6 种评测器", "✗ 仅支持精确断言", "◑ 依赖个人判断"],
    ["回归测试", "✓ 自动化版本对比", "✗ 不支持 AI 输出对比", "✗ 完全手动"],
    ["部署与数据", "✓ 本地优先 · 零后端 · 可移植", "◑ 需 CI/CD 集成", "✗ 无系统化方案"],
    ["多 AI 提供商", "✓ 3+ 提供商", "✗ N/A", "✗ N/A"],
]

table_left = Inches(1)
table_top = Inches(1.9)
rows = len(table_data)
cols = len(table_data[0])
col_widths = [Inches(1.8), Inches(3.2), Inches(3.2), Inches(3.2)]
row_height = Inches(0.7)

# Table header row
header_colors = [BORDER, ACCENT_BG, BORDER, BORDER]
for j in range(cols):
    x = table_left + sum(cw for cw in col_widths[:j])
    shape = add_rect(slide, x, table_top, col_widths[j], row_height,
                     fill_color=header_colors[j])
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = table_data[0][j]
    p.font.size = Pt(11)
    p.font.color.rgb = ACCENT_LIGHT if j == 1 else TEXT_MUTED
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

# Table body
for i in range(1, rows):
    for j in range(cols):
        x = table_left + sum(cw for cw in col_widths[:j])
        y = table_top + row_height + (i - 1) * row_height
        bg_color = CARD_BG if j != 1 or i == 0 else ACCENT_BG
        shape = add_rect(slide, x, y, col_widths[j], row_height,
                         fill_color=ACCENT_BG if j == 1 else CARD_BG,
                         border_color=BORDER)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = table_data[i][j]
        p.font.size = Pt(10)
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

        if j == 0:
            p.font.color.rgb = TEXT
            p.font.bold = True
        elif j == 1:
            p.font.color.rgb = GREEN if "✓" in table_data[i][j] else TEXT2
        elif "✗" in table_data[i][j]:
            p.font.color.rgb = RED
        elif "◑" in table_data[i][j]:
            p.font.color.rgb = YELLOW
        else:
            p.font.color.rgb = TEXT2

add_text(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.6),
         "一句话分析：AgentBench 是唯一专为 AI Agent 设计的评测运维平台。"
         "传统 QA 工具的断言模型无法处理 AI 输出的模糊性和多样性，手工评测则不可持续。",
         Pt(11), TEXT_MUTED)


# ═══════════════════════════════════════════
# SLIDE 7: 架构与技术栈
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "架构与技术栈", Pt(36), TEXT, bold=True)
add_text(slide, Inches(1), Inches(1.1), Inches(10), Inches(0.5),
         "极简架构，四个组件覆盖完整评测流程。", Pt(13), TEXT_MUTED)

arch_items = [
    ("💻", "前端", ACCENT_BG, ACCENT_LIGHT,
     [("Next.js 16 App Router", None, None, None),
      ("React 19 + TypeScript", None, None, None),
      ("Tailwind CSS 4 暗色主题", None, None, None),
      ("useSyncExternalStore", None, None, None)]),
    ("⚙", "API 层", GREEN_BG, GREEN,
     [("/api/execute — 代理转发", None, None, None),
      ("/api/analyze — AI 分析", None, None, None),
      ("密钥不落服务器", None, None, None),
      ("演示模式优雅降级", None, None, None)]),
    ("🛠", "评测引擎", YELLOW_BG, YELLOW,
     [("6 种评测器类型", None, None, None),
      ("LLM 评判含回退机制", None, None, None),
      ("Token 费用估算", None, None, None),
      ("可配置通过阈值", None, None, None)]),
    ("📦", "数据层", GREEN_BG, GREEN,
     [("套件 & 运行 → localStorage", None, None, None),
      ("设置 → localStorage", None, None, None),
      ("导出/导入 JSON + JSONL", None, None, None),
      ("首次访问自动播种", None, None, None)]),
]

for i, (icon, title, icon_bg, icon_clr, lines) in enumerate(arch_items):
    left = Inches(1 + i * 3.0)
    top = Inches(1.9)
    add_rect(slide, left, top, Inches(2.7), Inches(3.8))
    # Icon
    icon_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         left + Inches(0.25), top + Inches(0.25),
                                         Inches(0.5), Inches(0.5))
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = icon_bg
    icon_shape.line.fill.background()
    tf = icon_shape.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    # Title
    add_text(slide, left + Inches(0.9), top + Inches(0.25), Inches(1.5), Inches(0.4),
             title, Pt(15), icon_clr, bold=True)
    # Lines
    for j, (text, _, _, _) in enumerate(lines):
        add_text(slide, left + Inches(0.25), top + Inches(0.9 + j * 0.55),
                 Inches(2.2), Inches(0.4), text, Pt(10), TEXT_MUTED)

add_text(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
         "MIT 开源 · 无需后端 · 无需数据库 · 可自行部署",
         Pt(12), TEXT_MUTED, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SLIDE 8: 里程碑与时间轴
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "里程碑与发展计划", Pt(36), TEXT, bold=True)
add_text(slide, Inches(1), Inches(1.1), Inches(10), Inches(0.5),
         "已达成 5 个版本迭代，持续演进中。", Pt(13), TEXT_MUTED)

# Timeline items
timeline = [
    ("v0.1.0", "2026.05", "核心功能上线\n测试套件 + 运行引擎\n35 测试", True, False),
    ("v0.2.0", "2026.05", "CRUD 完善\n删除套件/运行\n确认弹窗", True, False),
    ("v0.3.0", "2026.05", "UX 交互优化\nEscape 关闭 · 点击外部\n共享 Hook 提取", True, False),
    ("v0.4.0", "2026.05", "测试扩展\n35→59 测试\n可访问性增强", True, False),
    ("v0.5.0", "2026.05 · 当前", "全面中文化\n代码评测器上线\nCC Switch 联动", True, True),
]

timeline_top = Inches(1.8)
timeline_y = timeline_top + Inches(0.3)
bar_left = Inches(1.2)
bar_right = Inches(12.1)
bar_width = bar_right - bar_left

# Timeline bar
add_rect(slide, bar_left, timeline_y, bar_width, Pt(3), fill_color=BORDER)

for i, (ver, date, desc, done, current) in enumerate(timeline):
    cx = bar_left + int(bar_width * (i / (len(timeline) - 1)))
    dot_color = ACCENT if current else (GREEN if done else BORDER)
    bar_color = GREEN if done else ACCENT if current else BORDER

    # Dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  cx - Inches(0.1), timeline_y - Inches(0.1),
                                  Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = dot_color
    dot.line.fill.background()

    # Version
    clr = ACCENT_LIGHT if current else TEXT
    add_text(slide, cx - Inches(0.9), timeline_y + Inches(0.3), Inches(1.8), Inches(0.3),
             ver, Pt(11), clr, bold=True, alignment=PP_ALIGN.CENTER, font_name=MONO)

    # Date
    add_text(slide, cx - Inches(0.9), timeline_y + Inches(0.55), Inches(1.8), Inches(0.3),
             date, Pt(8), TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # Description
    add_text(slide, cx - Inches(0.9), timeline_y + Inches(0.85), Inches(1.8), Inches(0.8),
             desc, Pt(9), TEXT2, alignment=PP_ALIGN.CENTER)

# Next steps
add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.4),
         "下一阶段计划", Pt(14), TEXT, bold=True)

next_steps = [
    "JSON Schema 全规范支持",
    "人工标注校准基准",
    "统计显著性检验",
    "CI/CD 集成",
]
for i, step in enumerate(next_steps):
    left = Inches(1 + i * 3.0)
    top = Inches(5.1)
    tag_shape = add_rect(slide, left, top, Inches(2.6), Inches(0.7), fill_color=CARD_BG, border_color=BORDER)
    tf = tag_shape.text_frame
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT2
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════
# SLIDE 9: 构建方式（团队法）
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "构建方式", Pt(36), TEXT, bold=True)
add_text(slide, Inches(1), Inches(1.1), Inches(10), Inches(0.5),
         "不是\"团队\"——是 AI 原生开发范式的一次验证。", Pt(13), TEXT_MUTED)

team_cards = [
    ("🧠", "100% AI 构建", ACCENT_BG, ACCENT_LIGHT,
     "整个 AgentBench 平台完全由 Claude Code 编写。从第一行代码到中文翻译到测试——每一行都是 AI 生成。",
     [("3", "轮迭代"), ("20+", "Bug 修复"), ("18+", "文件本地化")],
     ACCENT_LIGHT),
    ("🔗", "CC Switch 联动", GREEN_BG, GREEN,
     "通过 CC Switch 中间件采集真实 Claude Code 交互数据（5,340 次请求），自动生成编码任务测试套件，形成评测闭环。",
     [("5,340", "请求分析"), ("90%", "缓存命中"), ("4", "套件导出")],
     GREEN),
    ("🌐", "开源社区驱动", YELLOW_BG, YELLOW,
     "MIT 许可证，完全开源。任何 AI Agent 团队都可以自行部署、贡献测试套件、扩展评测器类型。",
     [("MIT", "许可证"), ("0", "外部依赖"), ("1", "命令部署")],
     YELLOW),
]

for i, (icon, title, icon_bg, icon_clr, desc, nums, num_clr) in enumerate(team_cards):
    left = Inches(1 + i * 3.9)
    top = Inches(1.9)
    add_rect(slide, left, top, Inches(3.6), Inches(4.5))
    # Icon
    icon_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         left + Inches(0.3), top + Inches(0.3),
                                         Inches(0.5), Inches(0.5))
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = icon_bg
    icon_shape.line.fill.background()
    tf = icon_shape.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    # Title
    add_text(slide, left + Inches(1.0), top + Inches(0.3), Inches(2.3), Inches(0.4),
             title, Pt(14), TEXT, bold=True)
    # Description
    add_text(slide, left + Inches(0.3), top + Inches(1.0), Inches(3.0), Inches(1.8),
             desc, Pt(10), TEXT2)
    # Three stat numbers
    for j, (val, label) in enumerate(nums):
        nx = left + Inches(0.3 + j * 1.15)
        ny = top + Inches(2.9)
        add_text(slide, nx, ny, Inches(1.0), Inches(0.4),
                 val, Pt(18), num_clr, bold=True, font_name=MONO)
        add_text(slide, nx, ny + Inches(0.35), Inches(1.0), Inches(0.3),
                 label, Pt(8), TEXT_MUTED)


# ═══════════════════════════════════════════
# SLIDE 10: 结尾 CTA
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)

add_badge(slide, Inches(5.6), Inches(1.2), "✓  可演示 · MIT 开源",
          bg_color=GREEN_BG, text_color=GREEN)

add_text(slide, Inches(1), Inches(1.9), Inches(11.3), Inches(1.2),
         "AgentBench", Pt(68), ACCENT_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(0.6),
         "AI Agent 评测运维平台 — 用 AI 测试 AI",
         Pt(20), TEXT2, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.5),
         "附创始人联系方式 · GitHub 二维码 · 开源地址",
         Pt(12), TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# CTA stats
cta_stats = [
    ("6", "评测器类型", ACCENT_LIGHT),
    ("72", "测试通过", ACCENT_LIGHT),
    ("3+", "AI 提供商", GREEN),
    ("0", "后端依赖", GREEN),
]
for i, (val, label, clr) in enumerate(cta_stats):
    left = Inches(2.5 + i * 2.3)
    top = Inches(4.5)
    add_rect(slide, left, top, Inches(1.9), Inches(1.5))
    add_text(slide, left, top + Inches(0.15), Inches(1.9), Inches(0.6),
             val, Pt(28), clr, bold=True, alignment=PP_ALIGN.CENTER, font_name=MONO)
    add_text(slide, left, top + Inches(0.85), Inches(1.9), Inches(0.3),
             label, Pt(10), TEXT_MUTED, alignment=PP_ALIGN.CENTER)

meta_items_cta = [
    "MIT 许可证",
    "github.com/LinguistWantsTech/agentbench",
    "由 Claude Code 构建",
]
for i, item in enumerate(meta_items_cta):
    add_text(slide, Inches(2.5 + i * 3.2), Inches(6.5), Inches(3.0), Inches(0.4),
             "●  " + item, Pt(11), TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# ── Save ──
output_path = os.path.join(os.path.dirname(__file__), "AgentBench-路演PPT.pptx")
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"Slides: {len(prs.slides)}")
